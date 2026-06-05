#!/usr/bin/env python3
"""Generate Financial-Manager-Presentation.pptx (16:9) for Google Slides import."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
EMERALD      = RGBColor(0x10, 0xB9, 0x81)
EMERALD_DARK = RGBColor(0x05, 0x96, 0x69)
EMERALD_SOFT = RGBColor(0xD1, 0xFA, 0xE5)
SLATE_900    = RGBColor(0x0F, 0x17, 0x2A)
SLATE_800    = RGBColor(0x1E, 0x29, 0x3B)
SLATE_700    = RGBColor(0x33, 0x41, 0x55)
SLATE_600    = RGBColor(0x47, 0x55, 0x69)
SLATE_500    = RGBColor(0x64, 0x74, 0x8B)
SLATE_400    = RGBColor(0x94, 0xA3, 0xB8)
SLATE_200    = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100    = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_50     = RGBColor(0xF8, 0xFA, 0xFC)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
AMBER        = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_SOFT   = RGBColor(0xFE, 0xF3, 0xC7)
BLUE         = RGBColor(0x3B, 0x82, 0xF6)
BLUE_SOFT    = RGBColor(0xDB, 0xEA, 0xFE)
VIOLET       = RGBColor(0x8B, 0x5C, 0xF6)
VIOLET_SOFT  = RGBColor(0xED, 0xE9, 0xFE)
ROSE         = RGBColor(0xF4, 0x3F, 0x5E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

FONT = "Arial"

# ---------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def _set_font(run, size, color, bold, italic, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font

def text(s, x, y, w, h, lines, size=18, color=SLATE_800, bold=False,
         italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         font=FONT, space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(ln, tuple):
            segs, lsize, lcolor, lbold = ln
        else:
            segs, lsize, lcolor, lbold = ln, size, color, bold
        if isinstance(segs, str):
            segs = [(segs, lcolor, lbold)]
        for seg in segs:
            txt, scolor, sbold = seg
            r = p.add_run()
            r.text = txt
            _set_font(r, lsize, scolor, sbold, italic, font)
    return tb

def box(s, x, y, w, h, fill=WHITE, line=None, line_w=1.0, shadow=False,
        radius=0.08, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {})
        sh = ef.makeelement(qn('a:outerShdw'),
                            {'blurRad': '90000', 'dist': '38100',
                             'dir': '5400000', 'rotWithShape': '0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val': '0F172A'})
        alpha = clr.makeelement(qn('a:alpha'), {'val': '22000'})
        clr.append(alpha); sh.append(clr); ef.append(sh); el.append(ef)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def boxtext(s, x, y, w, h, lines, fill=WHITE, line=None, line_w=1.0,
            size=14, color=SLATE_800, bold=False, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, shadow=False, radius=0.08,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, space_after=3, line_spacing=1.0):
    box(s, x, y, w, h, fill, line, line_w, shadow, radius, shape)
    pad = Inches(0.08)
    text(s, x + pad, y, w - 2 * pad, h, lines, size, color, bold,
         align=align, anchor=anchor, space_after=space_after, line_spacing=line_spacing)

def connector(s, x1, y1, x2, y2, color=SLATE_400, w=1.75, dash=None, arrow=True):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    ln = cn.line._get_or_add_ln()
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)
    if arrow:
        tail = ln.makeelement(qn('a:tailEnd'),
                              {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    return cn

def header(s, kicker, title, idx, total=9, dark=False):
    """Content-slide header with emerald accent + page number."""
    box(s, 0, 0, Inches(0.18), SH, fill=EMERALD)
    kc = EMERALD if not dark else EMERALD
    tc = SLATE_900 if not dark else WHITE
    text(s, Inches(0.6), Inches(0.34), Inches(11.5), Inches(0.4),
         kicker.upper(), 13, kc, True)
    text(s, Inches(0.6), Inches(0.66), Inches(11.5), Inches(0.7),
         title, 30, tc, True)
    text(s, Inches(12.2), Inches(0.42), Inches(0.9), Inches(0.4),
         f"{idx:02d} / {total:02d}", 11, SLATE_400, True, align=PP_ALIGN.RIGHT)
    box(s, Inches(0.6), Inches(1.42), Inches(11.0), Pt(2), fill=SLATE_200, radius=0)

def chip(s, x, y, w, label, fill, txtcolor, h=Inches(0.34), size=11):
    boxtext(s, x, y, w, h, label, fill=fill, color=txtcolor, bold=True,
            size=size, radius=0.5)

# ================================================================ SLIDE 1 — TITLE
s = slide()
bg(s, SLATE_900)
# subtle accent band
box(s, 0, Inches(6.95), SW, Inches(0.55), fill=EMERALD)
box(s, 0, Inches(6.85), SW, Inches(0.1), fill=EMERALD_DARK)
# logo mark
box(s, Inches(0.9), Inches(0.95), Inches(0.95), Inches(0.95),
    fill=EMERALD, radius=0.28)
text(s, Inches(0.9), Inches(0.95), Inches(0.95), Inches(0.95), "$",
     46, SLATE_900, True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(2.05), Inches(1.12), Inches(8), Inches(0.6),
     "FINANCIAL MANAGER", 15, EMERALD, True)
text(s, Inches(2.07), Inches(1.46), Inches(8), Inches(0.5),
     "Local-first personal finance, encrypted end-to-end", 13, SLATE_400, False)

text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.8),
     ["A private finance tracker", "with a built-in AI assistant"],
     50, WHITE, True, line_spacing=1.02)
text(s, Inches(0.92), Inches(4.55), Inches(11), Inches(0.8),
     "Your money data never leaves your machine unencrypted — and an AI assistant "
     "that respects that boundary by design.", 17, SLATE_400, line_spacing=1.2)

# feature pills
pills = [("AES-256-GCM encrypted at rest", EMERALD),
         ("Per-user SQLite", BLUE),
         ("Claude-powered chat", VIOLET),
         ("Self-hostable", AMBER)]
px = Inches(0.9)
for label, c in pills:
    w = Inches(0.42 + 0.105 * len(label))
    boxtext(s, px, Inches(5.55), w, Inches(0.46), label, fill=SLATE_800,
            line=c, line_w=1.25, color=WHITE, bold=True, size=11.5, radius=0.5)
    px = px + w + Inches(0.2)

text(s, Inches(0.9), Inches(7.0), Inches(8), Inches(0.45),
     "Team Presentation  ·  June 2026", 12.5, SLATE_900, True,
     anchor=MSO_ANCHOR.MIDDLE)

# ================================================================ SLIDE 2 — INTRODUCTION
s = slide()
bg(s, WHITE)
header(s, "Introduction", "Concept, Motive & Purpose", 2)

cards = [
    ("CONCEPT", EMERALD, EMERALD_SOFT,
     "What it is",
     ["A local-first personal finance tracker.",
      "Accounts, transactions, budgets, net worth,",
      "goals & credit — plus an AI assistant.",
      "No remote database: everything lives in",
      "per-user encrypted SQLite files on disk."]),
    ("MOTIVE", BLUE, BLUE_SOFT,
     "Why build it",
     ["Finance apps demand the most sensitive data",
      "you have — then store it on someone",
      "else's servers.",
      "We wanted full insight into our money",
      "without surrendering custody of it."]),
    ("PURPOSE", VIOLET, VIOLET_SOFT,
     "What it delivers",
     ["Bank-grade encryption on every value,",
      "kept entirely under user control.",
      "A genuinely useful AI assistant that",
      "only ever sees aggregates — unless you",
      "explicitly opt in, per message."]),
]
cw, gap = Inches(3.86), Inches(0.34)
cx = Inches(0.6)
cy = Inches(1.75)
ch = Inches(4.0)
for kicker, accent, soft, head, body in cards:
    box(s, cx, cy, cw, ch, fill=SLATE_50, line=SLATE_200, line_w=1.0,
        radius=0.05, shadow=True)
    box(s, cx, cy, cw, Inches(0.12), fill=accent, radius=0)
    chip(s, cx + Inches(0.3), cy + Inches(0.32), Inches(1.4), kicker, soft, accent)
    text(s, cx + Inches(0.3), cy + Inches(0.82), cw - Inches(0.6), Inches(0.5),
         head, 21, SLATE_900, True)
    text(s, cx + Inches(0.3), cy + Inches(1.4), cw - Inches(0.55), Inches(2.4),
         body, 13.5, SLATE_600, line_spacing=1.22, space_after=2)
    cx = cx + cw + gap

boxtext(s, Inches(0.6), Inches(6.05), Inches(12.13), Inches(0.95),
        [(([("Guiding principle:  ", SLATE_900, True),
            ("“If the server is compromised, the attacker still gets nothing but ciphertext.”",
             EMERALD_DARK, True)]), 16, SLATE_900, True)],
        fill=EMERALD_SOFT, line=EMERALD, line_w=1.25, anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER, radius=0.12)

# ================================================================ SLIDE 3 — LOCAL TECH STACK
s = slide()
bg(s, WHITE)
header(s, "Architecture", "Local Tech Stack", 3)

# left: layered request-flow diagram
dx = Inches(0.6)
dy = Inches(1.75)
dw = Inches(7.7)
layers = [
    ("Browser  ·  React 19 UI", "Client components · Tailwind CSS · Recharts · dnd-kit", BLUE, BLUE_SOFT),
    ("Next.js 15 App Router", "Pages in app/  ·  server + client rendering", EMERALD, EMERALD_SOFT),
    ("API Routes  ·  app/api/*", "Every route → getServerSession() → 401 if locked", EMERALD, EMERALD_SOFT),
    ("lib/  service layer", "crypto.ts · session.ts · server-session.ts · db.ts", VIOLET, VIOLET_SOFT),
    ("Per-user encrypted SQLite", "data/users/{id}/finance.db  ·  better-sqlite3 · WAL", AMBER, AMBER_SOFT),
]
lh = Inches(0.84)
lgap = Inches(0.18)
ly = dy
for i, (title, sub, accent, soft) in enumerate(layers):
    box(s, dx, ly, dw, lh, fill=soft, line=accent, line_w=1.25, radius=0.1)
    box(s, dx, ly, Inches(0.14), lh, fill=accent, radius=0)
    text(s, dx + Inches(0.35), ly + Inches(0.1), dw - Inches(0.6), Inches(0.4),
         title, 16, SLATE_900, True)
    text(s, dx + Inches(0.35), ly + Inches(0.46), dw - Inches(0.6), Inches(0.34),
         sub, 11.5, SLATE_600)
    if i < len(layers) - 1:
        mid = dx + dw / 2
        connector(s, mid, ly + lh, mid, ly + lh + lgap, color=SLATE_400, w=2.0)
    ly = ly + lh + lgap

# right: stack list + encryption note
rx = Inches(8.6)
rw = Inches(4.13)
box(s, rx, Inches(1.75), rw, Inches(3.05), fill=SLATE_900, radius=0.06, shadow=True)
text(s, rx + Inches(0.32), Inches(1.95), rw - Inches(0.6), Inches(0.4),
     "STACK AT A GLANCE", 12, EMERALD, True)
techlines = [
    ("Framework", "Next.js 15 · React 19"),
    ("Language", "TypeScript 5"),
    ("Styling", "Tailwind CSS 3 (dark mode)"),
    ("Database", "better-sqlite3 (per user)"),
    ("Charts", "Recharts"),
    ("Drag & drop", "@dnd-kit"),
    ("AI", "@anthropic-ai/sdk"),
]
ty = Inches(2.42)
for k, v in techlines:
    text(s, rx + Inches(0.32), ty, Inches(1.55), Inches(0.32), k, 11.5, SLATE_400, True)
    text(s, rx + Inches(1.7), ty, rw - Inches(2.0), Inches(0.32), v, 11.5, WHITE, True)
    ty = ty + Inches(0.323)

box(s, rx, Inches(4.95), rw, Inches(2.05), fill=EMERALD_SOFT, line=EMERALD,
    line_w=1.25, radius=0.07)
text(s, rx + Inches(0.32), Inches(5.12), rw - Inches(0.6), Inches(0.4),
     "🔒  ENCRYPTED AT REST", 12.5, EMERALD_DARK, True)
text(s, rx + Inches(0.32), Inches(5.5), rw - Inches(0.6), Inches(1.4),
     ["AES-256-GCM on every sensitive field (the _enc columns).",
      "Key = PBKDF2-SHA256, 310,000 iterations from the user's password.",
      "The key lives only in a signed cookie — never on disk."],
     12, SLATE_700, line_spacing=1.18, space_after=4)

# ================================================================ SLIDE 4 — DEPLOYED TECH STACK
s = slide()
bg(s, WHITE)
header(s, "Architecture", "Deployed Tech Stack", 4)

# horizontal pipeline diagram
dy = Inches(2.15)
bh = Inches(1.7)
steps = [
    ("🌐", "Internet", "HTTPS  ·  force_https", SLATE_700, SLATE_100),
    ("🎈", "Fly.io edge", "Anycast  ·  region iad\nauto-starts on request", BLUE, BLUE_SOFT),
    ("🖥️", "Firecracker microVM", "shared-cpu-1x · 512 MB\nscale-to-zero when idle", AMBER, AMBER_SOFT),
    ("▲", "Next.js standalone", "Docker image\nnode server.js  ·  :3000", EMERALD, EMERALD_SOFT),
    ("💾", "Fly Volume  fm_data", "3 GB → /app/data\nencrypted SQLite persists", VIOLET, VIOLET_SOFT),
]
n = len(steps)
bw = Inches(2.16)
gap = (dw_total := SW - Inches(1.2) - bw * n) / (n - 1)
x = Inches(0.6)
centers = []
for emoji, title, sub, accent, soft in steps:
    box(s, x, dy, bw, bh, fill=soft, line=accent, line_w=1.5, radius=0.1, shadow=True)
    text(s, x, dy + Inches(0.14), bw, Inches(0.5), emoji, 26, SLATE_900,
         align=PP_ALIGN.CENTER)
    text(s, x, dy + Inches(0.66), bw, Inches(0.4), title, 14.5, SLATE_900, True,
         align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.08), dy + Inches(1.02), bw - Inches(0.16), Inches(0.6),
         sub.split("\n"), 10.5, SLATE_600, align=PP_ALIGN.CENTER, line_spacing=1.1,
         space_after=0)
    centers.append((x, x + bw))
    x = x + bw + gap
for i in range(n - 1):
    connector(s, centers[i][1], dy + bh / 2, centers[i + 1][0], dy + bh / 2,
              color=SLATE_500, w=2.25)

# bootstrap + deploy strip
box(s, Inches(0.6), Inches(4.35), Inches(7.6), Inches(2.55), fill=SLATE_900,
    radius=0.06, shadow=True)
text(s, Inches(0.92), Inches(4.55), Inches(7), Inches(0.4),
     "HOW IT'S DEPLOYED", 12, EMERALD, True)
ti = [
    ("flyctl deploy", "Builds the multi-stage Dockerfile, pushes the image to Fly's registry, and boots a Firecracker microVM — no servers to manage."),
    ("Scale-to-zero", "auto_stop_machines + min_machines_running = 0: the VM sleeps when idle and auto-starts on the next request — pay ~nothing unused."),
    ("Multi-stage build", "Alpine build stage compiles the better-sqlite3 native module; runtime ships only the Next.js standalone output."),
]
yy = Inches(4.95)
for k, v in ti:
    text(s, Inches(0.92), yy, Inches(1.7), Inches(0.5), k, 12.5, EMERALD, True)
    text(s, Inches(2.65), yy, Inches(5.4), Inches(0.6), v, 11.5, SLATE_200,
         line_spacing=1.12)
    yy = yy + Inches(0.62)

box(s, Inches(8.45), Inches(4.35), Inches(4.28), Inches(2.55), fill=EMERALD_SOFT,
    line=EMERALD, line_w=1.25, radius=0.07)
text(s, Inches(8.75), Inches(4.55), Inches(3.7), Inches(0.4),
     "ONE-COMMAND DEPLOY", 12.5, EMERALD_DARK, True)
text(s, Inches(8.75), Inches(4.98), Inches(3.7), Inches(0.5),
     [(([("$ ", EMERALD_DARK, True), ("fly deploy", SLATE_900, True)]), 13, SLATE_900, True)],
     font="Courier New")
text(s, Inches(8.75), Inches(5.45), Inches(3.75), Inches(1.4),
     ["• Persistent Fly Volume (fm_data, 3 GB) keeps every per-user finance.db across restarts.",
      "• force_https + shared-cpu-1x microVM in iad (US-East).",
      "• Same Dockerfile runs locally via docker compose."],
     11.5, SLATE_700, line_spacing=1.18, space_after=4)

# ================================================================ SLIDE 5 — CLAUDE INTEGRATIONS
s = slide()
bg(s, WHITE)
header(s, "Developer Workflow", "Claude Code Integrations", 5)

text(s, Inches(0.6), Inches(1.55), Inches(12), Inches(0.5),
     [(([("Project-scoped ", SLATE_700, False),
         ("agents", VIOLET, True),
         (" and ", SLATE_700, False),
         ("skills", EMERALD_DARK, True),
         (" live in ", SLATE_700, False),
         (".claude/", SLATE_900, True),
         (" — they teach Claude this repo's exact conventions, so changes land consistent and review-ready.",
          SLATE_700, False)]), 14, SLATE_700, False)],
     line_spacing=1.1)

# Agents column
ax = Inches(0.6)
aw = Inches(5.95)
box(s, ax, Inches(2.25), aw, Inches(2.45), fill=VIOLET_SOFT, line=VIOLET,
    line_w=1.25, radius=0.05)
text(s, ax + Inches(0.3), Inches(2.4), aw - Inches(0.6), Inches(0.45),
     "🤖  AGENTS  ·  .claude/agents/", 14, VIOLET, True)
agents = [
    ("api-route", "Scaffolds an app/api route with the session guard, encrypt/decrypt rules, uuidv4 IDs & transaction wrapping baked in."),
    ("db-migration", "Writes a correct migrateToVN(), guards on user_version, and wires it into initSchema()."),
    ("feature-page", "Builds a new page with the auth gate, fetch pattern, dark-mode Tailwind & nav entry."),
]
yy = Inches(2.85)
for k, v in agents:
    text(s, ax + Inches(0.3), yy, Inches(1.55), Inches(0.5), k, 12.5, VIOLET, True)
    text(s, ax + Inches(1.85), yy, aw - Inches(2.15), Inches(0.65), v, 11, SLATE_700,
         line_spacing=1.1)
    yy = yy + Inches(0.6)

# Skills column
sx = Inches(6.78)
sw = Inches(5.95)
box(s, sx, Inches(2.25), sw, Inches(2.45), fill=EMERALD_SOFT, line=EMERALD,
    line_w=1.25, radius=0.05)
text(s, sx + Inches(0.3), Inches(2.4), sw - Inches(0.6), Inches(0.45),
     "📚  SKILLS  ·  .claude/skills/", 14, EMERALD_DARK, True)
skills = [
    ("add-api-route", "Full GET/POST/PUT/DELETE template + the route's hard rules."),
    ("add-migration", "Step-by-step schema-change recipe with field-naming conventions."),
    ("add-page", "Auth gating, loading skeletons, nav wiring & dark-mode pairs."),
    ("encryption-patterns", "Reference for deriveKey / encrypt / decrypt & nullable _enc fields."),
]
yy = Inches(2.82)
for k, v in skills:
    text(s, sx + Inches(0.3), yy, Inches(1.85), Inches(0.5), k, 12, EMERALD_DARK, True)
    text(s, sx + Inches(2.15), yy, sw - Inches(2.45), Inches(0.5), v, 11, SLATE_700,
         line_spacing=1.08)
    yy = yy + Inches(0.46)

# why it helps strip
box(s, Inches(0.6), Inches(4.95), Inches(12.13), Inches(2.0), fill=SLATE_900,
    radius=0.05, shadow=True)
text(s, Inches(0.92), Inches(5.15), Inches(11), Inches(0.4),
     "WHY IT SPEEDS UP DEVELOPMENT", 12.5, EMERALD, True)
bens = [
    ("Conventions, encoded", "Encryption rules, the session guard & migration pattern are captured once — not re-explained each task."),
    ("Consistency by default", "Every new route/page/migration follows the same shape, so reviews stay fast and bugs stay rare."),
    ("Security stays intact", "Agents never emit a plaintext column for an _enc field or skip getServerSession() — the risky mistakes."),
]
bx = Inches(0.92)
bw2 = Inches(3.78)
for k, v in bens:
    text(s, bx, Inches(5.6), bw2, Inches(0.45), k, 14, EMERALD, True)
    text(s, bx, Inches(6.02), bw2, Inches(0.9), v, 11.5, SLATE_200, line_spacing=1.15)
    bx = bx + bw2 + Inches(0.39)

# ================================================================ SLIDE 6 — APP FEATURES
s = slide()
bg(s, WHITE)
header(s, "Product Tour", "App Features", 6)

feats = [
    ("📊", "Dashboard", "Monthly income/expense, savings rate, category donut, 6-month trends, budget & debt progress, quick-log.", EMERALD),
    ("🏦", "Accounts", "Drag-and-drop list across 6 account types; balances, APR, limits; statement modal with running balance.", BLUE),
    ("💳", "Transactions", "Monthly filtered ledger; add/edit/delete; income/expense/transfer; recurring templates.", VIOLET),
    ("🧾", "Spending", "Category-level analysis vs. prior periods, top-categories ranking, budget status.", AMBER),
    ("📈", "Wealth", "Net worth over time, asset-vs-debt donut, account composition, debt-to-asset ratio.", EMERALD),
    ("🗺️", "Planning", "Cash-flow forecast of scheduled transactions + a goals board with target dates.", BLUE),
    ("⭐", "Credit", "Credit-score history, per-card utilization (30% marker), debt-payoff progress.", ROSE),
    ("🎯", "Goals", "Color-coded savings goals with targets, dates & progress visualization.", VIOLET),
    ("📄", "Reports", "Annual summaries, CSV export & charts for deeper review.", AMBER),
    ("👤", "Profile", "Username/password change, encrypted data export, account settings.", SLATE_600),
    ("💬", "AI Chat", "Streaming finance assistant with encrypted session history (next slide).", EMERALD_DARK),
    ("🌙", "Dark mode", "Full dark-mode theme toggled across every page.", SLATE_700),
]
cols, rows = 4, 3
gx, gy = Inches(0.6), Inches(1.7)
cw = Inches(2.95)
chh = Inches(1.62)
hgap = Inches(0.083)
vgap = Inches(0.13)
for i, (emoji, title, body, accent) in enumerate(feats):
    r, c = divmod(i, cols)
    x = gx + c * (cw + hgap)
    y = gy + r * (chh + vgap)
    box(s, x, y, cw, chh, fill=SLATE_50, line=SLATE_200, line_w=1.0, radius=0.08)
    box(s, x, y, Inches(0.1), chh, fill=accent, radius=0)
    text(s, x + Inches(0.25), y + Inches(0.14), Inches(0.6), Inches(0.45),
         emoji, 19)
    text(s, x + Inches(0.82), y + Inches(0.17), cw - Inches(0.95), Inches(0.4),
         title, 14.5, SLATE_900, True)
    text(s, x + Inches(0.25), y + Inches(0.66), cw - Inches(0.45), Inches(0.9),
         body, 10.3, SLATE_600, line_spacing=1.08)

# ================================================================ SLIDE 7 — AI CHAT PIPELINE
s = slide()
bg(s, WHITE)
header(s, "AI Assistant", "AI Chat — How It Works", 7)

text(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.45),
     [(([("A streamlined ", SLATE_700, False),
         ("2-call pipeline", VIOLET, True),
         (" — the only place decrypted data leaves the device, and even then only aggregates by default.",
          SLATE_700, False)]), 14, SLATE_700, False)])

# pipeline diagram
dy = Inches(2.25)
bh = Inches(1.55)
# user box
box(s, Inches(0.6), dy + Inches(0.25), Inches(1.7), Inches(1.05), fill=SLATE_100,
    line=SLATE_300 if False else SLATE_400, line_w=1.25, radius=0.12)
text(s, Inches(0.6), dy + Inches(0.25), Inches(1.7), Inches(1.05),
     ["👤", "User message"], 17, SLATE_800, True, align=PP_ALIGN.CENTER,
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)

# step 1
s1x = Inches(2.75)
sw1 = Inches(3.35)
box(s, s1x, dy, sw1, bh, fill=BLUE_SOFT, line=BLUE, line_w=1.5, radius=0.08, shadow=True)
chip(s, s1x + Inches(0.25), dy + Inches(0.18), Inches(1.0), "CALL 1", BLUE, WHITE, size=10)
chip(s, s1x + Inches(1.32), dy + Inches(0.18), Inches(1.78), "claude-haiku-4-5", WHITE, BLUE, size=9.5)
text(s, s1x + Inches(0.25), dy + Inches(0.6), sw1 - Inches(0.5), Inches(0.4),
     "classifyAndGuard()", 15, SLATE_900, True)
text(s, s1x + Inches(0.25), dy + Inches(0.98), sw1 - Inches(0.5), Inches(0.5),
     "Gate unsafe / off-topic → pick lookups  ·  cheap, fast guard", 10.5,
     SLATE_700, line_spacing=1.05)

# step 2
s2x = Inches(6.55)
sw2 = Inches(3.35)
box(s, s2x, dy, sw2, bh, fill=VIOLET_SOFT, line=VIOLET, line_w=1.5, radius=0.08, shadow=True)
chip(s, s2x + Inches(0.25), dy + Inches(0.18), Inches(1.0), "CALL 2", VIOLET, WHITE, size=10)
chip(s, s2x + Inches(1.32), dy + Inches(0.18), Inches(1.85), "claude-sonnet-4-6", WHITE, VIOLET, size=9.5)
text(s, s2x + Inches(0.25), dy + Inches(0.6), sw2 - Inches(0.5), Inches(0.4),
     "streamAnswer()", 15, SLATE_900, True)
text(s, s2x + Inches(0.25), dy + Inches(0.98), sw2 - Inches(0.5), Inches(0.5),
     "Cached system + aggregates → streams the reply", 10.5,
     SLATE_700, line_spacing=1.05)

# response box
rxb = Inches(10.35)
box(s, rxb, dy + Inches(0.25), Inches(2.35), Inches(1.05), fill=EMERALD_SOFT,
    line=EMERALD, line_w=1.5, radius=0.12)
text(s, rxb, dy + Inches(0.25), Inches(2.35), Inches(1.05),
     ["💬 SSE stream", "deltas → suggestions", "→ done"], 11.5, SLATE_900, True,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

# arrows
connector(s, Inches(2.3), dy + bh / 2, s1x, dy + bh / 2, color=SLATE_500, w=2.25)
connector(s, s1x + sw1, dy + bh / 2, s2x, dy + bh / 2, color=SLATE_500, w=2.25)
connector(s, s2x + sw2, dy + bh / 2, rxb, dy + bh / 2, color=SLATE_500, w=2.25)
text(s, s1x + sw1, dy - Inches(0.05), Inches(0.45), Inches(0.5), "intents",
     8.5, SLATE_500, True, align=PP_ALIGN.CENTER)

# model note
boxtext(s, Inches(2.75), dy + bh + Inches(0.22), Inches(7.15), Inches(0.5),
        [(([("Two-model split:  ", SLATE_700, False),
            ("Haiku 4.5", BLUE, True),
            (" guards & routes, ", SLATE_700, False),
            ("Sonnet 4.6", VIOLET, True),
            (" answers  ·  system prompt cached", SLATE_700, False)]), 12, SLATE_700, False)],
        fill=SLATE_50, line=SLATE_200, line_w=1.0, radius=0.2, anchor=MSO_ANCHOR.MIDDLE)

# privacy + persistence cards
py = Inches(4.7)
box(s, Inches(0.6), py, Inches(5.95), Inches(2.25), fill=EMERALD_SOFT, line=EMERALD,
    line_w=1.25, radius=0.05)
text(s, Inches(0.9), py + Inches(0.18), Inches(5.4), Inches(0.4),
     "🔐  PRIVACY MODEL", 13, EMERALD_DARK, True)
text(s, Inches(0.9), py + Inches(0.62), Inches(5.4), Inches(1.6),
     ["• Default: aggregates only — balances, category totals, budget status, net worth.",
      "• Row-level transaction detail is sent ONLY when you flip the per-message “Include transaction detail” toggle.",
      "• Nothing else ever leaves the device in plaintext."],
     12, SLATE_700, line_spacing=1.18, space_after=5)

box(s, Inches(6.78), py, Inches(5.95), Inches(2.25), fill=SLATE_900, radius=0.05, shadow=True)
text(s, Inches(7.08), py + Inches(0.18), Inches(5.4), Inches(0.4),
     "💾  ENCRYPTED CHAT HISTORY", 13, EMERALD, True)
text(s, Inches(7.08), py + Inches(0.62), Inches(5.4), Inches(1.6),
     ["• Sessions & messages persist in the per-user finance.db (migration v15).",
      "• Stored as title_enc / content_enc — same AES-256-GCM key as the rest of your data.",
      "• Multi-turn history, rename & delete sessions; suggested prompts to start."],
     12, SLATE_200, line_spacing=1.18, space_after=5)

# ================================================================ SLIDE 8 — AI IMPLEMENTATION DETAIL
s = slide()
bg(s, WHITE)
header(s, "AI Assistant", "Implementation Details", 8)

# code-ish module map
box(s, Inches(0.6), Inches(1.7), Inches(5.95), Inches(3.05), fill=SLATE_900,
    radius=0.05, shadow=True)
text(s, Inches(0.92), Inches(1.9), Inches(5.4), Inches(0.4),
     "lib/ai/  —  the assistant in four files", 13.5, EMERALD, True)
mods = [
    ("client.ts", "Anthropic client + model constants"),
    ("pipeline.ts", "classifyAndGuard() + streamAnswer()"),
    ("finance-context.ts", "Decrypts DB → aggregate context"),
    ("chat-store.ts", "Encrypted sessions / messages CRUD"),
]
yy = Inches(2.42)
for k, v in mods:
    text(s, Inches(0.92), yy, Inches(2.3), Inches(0.4),
         k, 13, EMERALD, True, font="Courier New")
    text(s, Inches(3.0), yy, Inches(3.4), Inches(0.5), v, 11.5, SLATE_200,
         line_spacing=1.05)
    yy = yy + Inches(0.55)
text(s, Inches(0.92), Inches(4.5), Inches(5.4), Inches(0.3),
     "Streamed over SSE from app/api/chat/stream/route.ts", 10.5, SLATE_400, True)

# highlights — left-accent cards, no icons
hx = Inches(6.78)
hw = Inches(5.95)
hl = [
    ("Token streaming", "Server-Sent Events push session → deltas → suggestions → done; the UI renders the reply as it generates.", BLUE),
    ("Intent routing", "Call 1 returns structured intents (accounts, spending, income, budget, transactions) so only needed data is decrypted.", EMERALD),
    ("Guardrail first", "Unsafe or off-topic asks are refused before any financial data is touched.", AMBER),
    ("Same-key encryption", "Chat content reuses the user's derived key — no separate secret, no plaintext history.", VIOLET),
]
yy = Inches(1.7)
for k, v, accent in hl:
    box(s, hx, yy, hw, Inches(0.72), fill=SLATE_50, line=SLATE_200, line_w=1.0, radius=0.1)
    box(s, hx, yy, Inches(0.1), Inches(0.72), fill=accent, radius=0)
    text(s, hx + Inches(0.36), yy + Inches(0.1), hw - Inches(0.6), Inches(0.32),
         k, 13.5, SLATE_900, True)
    text(s, hx + Inches(0.36), yy + Inches(0.39), hw - Inches(0.6), Inches(0.32),
         v, 10.5, SLATE_600, line_spacing=1.08)
    yy = yy + Inches(0.79)

# recent optimizations strip — accent rules, no icons
oy = Inches(5.05)
box(s, Inches(0.6), oy, Inches(12.13), Inches(1.55), fill=SLATE_900, radius=0.05, shadow=True)
text(s, Inches(0.92), oy + Inches(0.18), Inches(11), Inches(0.4),
     "RECENT OPTIMIZATIONS", 12.5, EMERALD, True)
opts = [
    ("Haiku for the guard", BLUE,
     "Call 1 now runs on claude-haiku-4-5 — the routing & safety pass is ~3–5× cheaper and faster, with no quality loss."),
    ("Prompt caching", EMERALD,
     "The generator's system block is cached (ephemeral) — up to ~90% off repeated input tokens on follow-up turns."),
    ("Follow-up chips", VIOLET,
     "Suggestions generated on the cheap model after each answer streams — richer UX, zero added latency."),
]
ox = Inches(0.92)
ow = Inches(3.6)
ogap = Inches(0.3)
for k, accent, v in opts:
    box(s, ox, oy + Inches(0.62), Inches(0.06), Inches(0.78), fill=accent, radius=0)
    text(s, ox + Inches(0.22), oy + Inches(0.56), ow - Inches(0.22), Inches(0.34),
         k, 14, WHITE, True)
    text(s, ox + Inches(0.22), oy + Inches(0.92), ow - Inches(0.3), Inches(0.55),
         v, 10.5, SLATE_400, line_spacing=1.15)
    ox = ox + ow + ogap

text(s, Inches(0.6), Inches(6.78), Inches(12.13), Inches(0.45),
     [(([("Useful AI, without the privacy trade-off", EMERALD_DARK, True),
         ("    encrypted at rest  ·  self-hostable  ·  aggregates by default", SLATE_500, False)]), 13.5, SLATE_700, False)],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- save
out = "Financial-Manager-Presentation.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
